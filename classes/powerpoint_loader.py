from pptx import Presentation
from langchain.document_loaders.base import BaseLoader
from langchain.schema import Document

class PowerPointLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> list[Document]:
        presentation = Presentation(self.file_path)
        print(presentation)
        documents = []

        for slide_number, slide in enumerate(presentation.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

            document = Document(
                page_content=slide_text,
                metadata={"slide_number": slide_number}
            )
            documents.append(document)

        return documents

if __name__ == '__main__':
    # 使用例
    loader = PowerPointLoader("/data/会津若松市_ICTを活用した地域の魅力と健康増進事業_構築/10.基本設計/作業用/プレゼンテーション1.pptx")
    documents = loader.load()

    for doc in documents:
        print(doc.page_content)